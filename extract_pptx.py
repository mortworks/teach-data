#!/usr/bin/env python3
import sys
import os
import re
import hashlib
from pptx import Presentation


def slugify(text: str) -> str:
    """Make a filesystem-friendly slug from text."""
    return re.sub(r'[^a-z0-9]+', '-', text.lower()).strip('-')


def classify_slide(title: str) -> str:
    """Classify a slide by its title."""
    if not title:
        return "Notes"
    title_lower = title.lower()
    if "question" in title_lower or title_lower.endswith("?"):
        return "Question"
    if "key term" in title_lower or "keywords" in title_lower:
        return "Key Terms"
    return "Notes"


def extract_pptx(pptx_file: str, out_root: str):
    prs = Presentation(pptx_file)
    unit_name = os.path.splitext(os.path.basename(pptx_file))[0]
    unit_slug = slugify(unit_name)
    unit_dir = os.path.join(out_root, unit_slug)
    img_dir = os.path.join(unit_dir, "images")

    os.makedirs(img_dir, exist_ok=True)

    md_path = os.path.join(unit_dir, f"{unit_slug}.md")
    with open(md_path, "w", encoding="utf-8") as md:
        # YAML front matter
        md.write(f"---\ntopic: {unit_name}\n---\n\n")

        for i, slide in enumerate(prs.slides, start=1):
            # Get title (if available)
            title_shape = slide.shapes.title if slide.shapes.title else None
            title_text = title_shape.text.strip() if title_shape and title_shape.text else ""
            section = classify_slide(title_text)

            md.write(f"## {section}\n\n")

            # Collect text content
            text_items = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape != title_shape:  # avoid duplicating the title
                        text_items.append(shape.text.strip())

                # Extract images
                if shape.shape_type == 13:  # picture
                    image = shape.image
                    ext = image.ext
                    h = hashlib.sha1(image.blob).hexdigest()[:8]
                    img_filename = f"slide{i}-{h}.{ext}"
                    with open(os.path.join(img_dir, img_filename), "wb") as f:
                        f.write(image.blob)
                    md.write(f"![Slide {i}]({{< relref \"images/{img_filename}\" >}})\n\n")

            if text_items:
                for line in text_items:
                    # Add bullet point style for notes and terms
                    if section in ["Notes", "Key Terms"]:
                        md.write(f"- {line}\n")
                    else:
                        md.write(f"{line}\n")
                md.write("\n")

    print(f"✅ Extracted '{pptx_file}' → {unit_dir}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx.py path/to/file.pptx")
        sys.exit(1)

    pptx_file = sys.argv[1]
    out_root = "data/computerscience/units/gcse/ocr-j277"

    extract_pptx(pptx_file, out_root)
